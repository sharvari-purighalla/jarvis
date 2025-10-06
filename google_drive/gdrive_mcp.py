from __future__ import print_function
import os, io, mimetypes
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from PyPDF2 import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation

from mcp import tool
#from mcp.server.fastapi import FastAPIServer
import uvicorn

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']


# authentication
def authenticate():
    creds = None
    if os.path.exists('token.json'):
        creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file('credentials.json', SCOPES)
            creds = flow.run_local_server(port=0)
        with open('token.json', 'w') as token:
            token.write(creds.to_json())
    return creds


# listing files from gdrive
@server.tool()
def list_files(service, limit=10):
    results = service.files().list(
        pageSize=limit, fields="files(id, name, mimeType)"
    ).execute()
    items = results.get('files', [])
    for i, item in enumerate(items, start=1):
        print(f"{i}. {item['name']} ({item['mimeType']}) — {item['id']}")
    return items


# downloading and reading files from gdrive
@server.tool()
def download_and_read(service, file_id, mime_type):
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while done is False:
        status, done = downloader.next_chunk()

    fh.seek(0)
    text = extract_text_from_bytes(fh, mime_type)
    return text

# extracting text from different files types
def extract_text_from_bytes(fh, mime_type):
    if "pdf" in mime_type:
        reader = PdfReader(fh)
        return "\n".join([p.extract_text() or "" for p in reader.pages])

    elif "word" in mime_type or "docx" in mime_type:
        doc = Document(fh)
        return "\n".join([p.text for p in doc.paragraphs])

    elif "sheet" in mime_type or "excel" in mime_type:
        wb = openpyxl.load_workbook(fh)
        data = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                data.append("\t".join(str(c) for c in row if c))
        return "\n".join(data)

    elif "presentation" in mime_type or "powerpoint" in mime_type:
        prs = Presentation(fh)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    else:
        try:
            return fh.read().decode("utf-8", errors="ignore")
        except:
            return "[Unsupported file type or unreadable content]"



if __name__ == '__main__':
    creds = authenticate()
    service = build('drive', 'v3', credentials=creds)

    files = list_files(service, limit=5)
    choice = int(input("Enter file number to read: ")) - 1
    file = files[choice]

    text = download_and_read(service, file['id'], file['mimeType'])